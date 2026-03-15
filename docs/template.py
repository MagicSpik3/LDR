from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_template():
    prs = Presentation()

    # 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Use the first slide layout as the template layout
    layout = prs.slide_layouts[0]
    shapes = layout.shapes

    # 1. Corporate Blue Top Bar
    top_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        Inches(0.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(31, 78, 121)
    top_bar.line.fill.background()

    # 2. Footer Bar
    bottom_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.4),
        prs.slide_width,
        Inches(0.4)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(31, 78, 121)
    bottom_bar.line.fill.background()

    # 3. Logo Text
    logo_box = shapes.add_textbox(
        prs.slide_width - Inches(2.2),
        0,
        Inches(2),
        Inches(0.5)
    )

    p = logo_box.text_frame.paragraphs[0]
    p.text = "WAS RADAR"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    prs.save("template.pptx")
    print("template.pptx created successfully.")

if __name__ == "__main__":
    create_template()